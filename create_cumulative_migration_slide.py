#!/usr/bin/env python3
"""Editable Google Slides deck: cumulative migration progress with clickable navigation."""

import os
import subprocess
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from migration_actual_vs_planned_data import (
    ACTUAL_WEEKLY,
    FINAL_TARGET,
    CURRENT_AS_OF,
    CURRENT_CUMULATIVE,
    MIGRATION_PHASES,
    build_chart_rows,
    weekly_total,
)

COLOR_LITE = RGBColor(135, 206, 235)
COLOR_PRO = RGBColor(27, 79, 138)
COLOR_ENT = RGBColor(30, 41, 59)
COLOR_ACTUAL = RGBColor(46, 158, 91)
COLOR_PLANNED = RGBColor(27, 79, 138)
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(209, 213, 219)
COLOR_MID = RGBColor(100, 116, 139)
COLOR_WARN = RGBColor(220, 38, 38)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
CHART_PNG = "/workspace/migration_actual_vs_planned_chart.png"


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_textbox(slide, left, top, width, height, text, size=12, bold=False,
                color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_nav_bar(slide, slides_map, active_key, google_compatible=False):
    """Navigation bar — hyperlinks omitted in google_compatible mode (breaks Slides import)."""
    items = [
        ("overview", "Overview"),
        ("chart", "Chart"),
        ("data", "Cumulative Data"),
        ("weekly", "Weekly Detail"),
        ("phases", "Phases"),
    ]
    bar_top = Inches(6.75)
    btn_w = Inches(2.35)
    btn_h = Inches(0.45)
    gap = Inches(0.15)
    start_x = Inches(0.6)

    for i, (key, label) in enumerate(items):
        left = start_x + i * (btn_w + gap)
        is_active = key == active_key
        fill = COLOR_PLANNED if is_active else COLOR_WHITE
        text_color = COLOR_WHITE if is_active else COLOR_PLANNED
        btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, bar_top, btn_w, btn_h)
        set_fill(btn, fill)
        btn.line.color.rgb = COLOR_PLANNED
        btn.line.width = Pt(1)
        btn.adjustments[0] = 0.3
        if not google_compatible and not is_active:
            btn.click_action.target_slide = slides_map[key]
        tf = btn.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def style_header_cell(cell, text, bg=COLOR_PLANNED):
    cell.text = text
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.size = Pt(9)
    cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg


def style_cum_header_cell(cell, text):
    style_header_cell(cell, text, bg=COLOR_ACTUAL)


def build_overview_slide(slide, slides_map, rows, google_compatible=False):
    add_bg(slide)
    latest = [r for r in rows if r["cumulative_actual"] != ""][-1]
    planned = latest["planned_cumulative"] if latest["planned_cumulative"] != "" else "—"
    variance = latest["variance"] if latest["variance"] != "" else 0
    var_color = COLOR_ACTUAL if variance >= 0 else COLOR_WARN

    add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.7),
                "Room Migration Progress — Cumulative Totals", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.4),
                f"As of {latest['week']}  ·  Click any tab below to navigate", size=12,
                color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)

    cards = [
        ("Cumulative Total", f"{latest['cumulative_actual']:,}", "All weeks combined through this date", COLOR_ACTUAL),
        ("Planned Target", f"{planned:,}" if planned != "—" else "—", f"Plan for {latest['week']}", COLOR_PLANNED),
        ("Variance", f"{variance:+,}" if variance != "" else "—", "Actual minus planned", var_color),
        ("Goal", f"{FINAL_TARGET:,}", "Near-term account target", COLOR_PRO),
    ]
    card_w = Inches(2.85)
    card_h = Inches(1.65)
    gap = Inches(0.3)
    start_x = (SLIDE_W - (card_w * 4 + gap * 3)) / 2

    for i, (title, value, sub, color) in enumerate(cards):
        left = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.85), card_w, card_h)
        set_fill(card, COLOR_WHITE)
        card.line.color.rgb = COLOR_BORDER
        card.adjustments[0] = 0.08
        if not google_compatible:
            card.click_action.target_slide = slides_map["data"]
        add_textbox(slide, left + Inches(0.2), Inches(2.0), card_w - Inches(0.4), Inches(0.35),
                    title, size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(2.45), card_w - Inches(0.4), Inches(0.55),
                    value, size=26, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), Inches(3.05), card_w - Inches(0.4), Inches(0.35),
                    sub, size=9, color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)

    # Quick cumulative table (last 4 weeks)
    add_textbox(slide, Inches(0.8), Inches(3.85), Inches(5), Inches(0.35),
                "Recent cumulative totals (click Chart for full view →)", size=13, bold=True)
    actual_rows = [r for r in rows if r["cumulative_actual"] != ""][-4:]
    mini = slide.shapes.add_table(len(actual_rows) + 1, 4, Inches(0.8), Inches(4.25), Inches(7.5), Inches(1.8))
    table = mini.table
    for col, h in enumerate(["Week", "Weekly", "Cumulative Total", "Planned"]):
        style_header_cell(table.cell(0, col), h)
    for ri, row in enumerate(actual_rows, start=1):
        table.cell(ri, 0).text = row["week"]
        table.cell(ri, 1).text = f"{row['weekly_actual']:,}"
        table.cell(ri, 2).text = f"{row['cumulative_actual']:,}"
        table.cell(ri, 2).text_frame.paragraphs[0].font.bold = True
        table.cell(ri, 2).fill.solid()
        table.cell(ri, 2).fill.fore_color.rgb = RGBColor(232, 245, 237)
        p = row["planned_cumulative"]
        table.cell(ri, 3).text = f"{p:,}" if p != "" else "—"
        for c in range(4):
            table.cell(ri, c).text_frame.paragraphs[0].font.size = Pt(10)

    # Chart preview thumbnail — clickable
    if os.path.exists(CHART_PNG):
        pic = slide.shapes.add_picture(CHART_PNG, Inches(8.6), Inches(3.85), width=Inches(4.0))
        if not google_compatible:
            pic.click_action.target_slide = slides_map["chart"]
        add_textbox(slide, Inches(8.6), Inches(6.35), Inches(4.0), Inches(0.3),
                    "See Chart slide for full view", size=10, color=COLOR_PLANNED, align=PP_ALIGN.CENTER, bold=True)

    add_nav_bar(slide, slides_map, "overview", google_compatible)


def build_chart_slide(slide, slides_map, google_compatible=False):
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6),
                "Actual vs Planned — Cumulative Line Chart", size=24, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35),
                "Green = cumulative actual total per date  ·  Blue dashed = planned  ·  Click Data tab to edit numbers",
                size=11, color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)

    if os.path.exists(CHART_PNG):
        slide.shapes.add_picture(CHART_PNG, Inches(0.9), Inches(1.4), width=Inches(11.5))
    else:
        add_textbox(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                    "Chart image not found. Run: python3 create_actual_vs_planned_chart.py",
                    size=14, color=COLOR_WARN, align=PP_ALIGN.CENTER)

    if not google_compatible:
        hint = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(6.2), Inches(4.3), Inches(0.42))
        set_fill(hint, COLOR_ACTUAL)
        set_no_line(hint)
        hint.adjustments[0] = 0.4
        hint.click_action.target_slide = slides_map["data"]
        tf = hint.text_frame
        tf.paragraphs[0].text = "✎  Edit data on Cumulative Data slide"
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        add_textbox(slide, Inches(4.5), Inches(6.2), Inches(4.3), Inches(0.42),
                    "Edit numbers on slide 3: Cumulative Data", size=11, bold=True,
                    color=COLOR_ACTUAL, align=PP_ALIGN.CENTER)

    add_nav_bar(slide, slides_map, "chart", google_compatible)


def build_data_slide(slide, slides_map, rows, google_compatible=False):
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.55),
                "Cumulative Data (Editable)", size=24, bold=True)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35),
                "Edit Weekly columns — Cumulative Total = sum of all prior weeks through that date",
                size=11, color=COLOR_MID, italic=True)

    actual_rows = [r for r in rows if r["cumulative_actual"] != ""]
    future_rows = [r for r in rows if r["cumulative_actual"] == "" and r["planned_cumulative"] != ""]
    all_display = actual_rows + future_rows

    cols = ["Week", "Wk Lite", "Wk Pro", "Wk Ent", "Wk Total",
            "Cum Lite", "Cum Pro", "Cum Ent", "Cum Total", "Planned", "Variance"]
    table_shape = slide.shapes.add_table(
        len(all_display) + 1, len(cols),
        Inches(0.35), Inches(1.35), Inches(12.6), Inches(5.2),
    )
    table = table_shape.table

    for col, h in enumerate(cols):
        cell = table.cell(0, col)
        if h.startswith("Cum"):
            style_cum_header_cell(cell, h)
        else:
            style_header_cell(cell, h)

    for ri, row in enumerate(all_display, start=1):
        vals = [
            row["week"],
            str(row["weekly_lite"]) if row["weekly_lite"] != "" else "",
            str(row["weekly_pro"]) if row["weekly_pro"] != "" else "",
            str(row["weekly_enterprise"]) if row["weekly_enterprise"] != "" else "",
            str(row["weekly_actual"]) if row["weekly_actual"] != "" else "",
            str(row["cumulative_lite"]) if row["cumulative_lite"] != "" else "",
            str(row["cumulative_pro"]) if row["cumulative_pro"] != "" else "",
            str(row["cumulative_enterprise"]) if row["cumulative_enterprise"] != "" else "",
            str(row["cumulative_actual"]) if row["cumulative_actual"] != "" else "",
            str(row["planned_cumulative"]) if row["planned_cumulative"] != "" else "",
            f"{row['variance']:+,}" if row["variance"] != "" else "",
        ]
        for ci, val in enumerate(vals):
            cell = table.cell(ri, ci)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(8)
            if ci == 8 and val:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(232, 245, 237)

    add_nav_bar(slide, slides_map, "data", google_compatible)


def build_weekly_slide(slide, slides_map, google_compatible=False):
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.55),
                "Weekly Detail by Segment (Editable)", size=24, bold=True)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35),
                "Enter weekly counts only — not cumulative. Totals roll up on Cumulative Data slide.",
                size=11, color=COLOR_MID, italic=True)

    headers = ["Week Starting", "Lite", "Pro", "Enterprise", "Weekly Total"]
    table_shape = slide.shapes.add_table(
        len(ACTUAL_WEEKLY) + 1, len(headers),
        Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.8),
    )
    table = table_shape.table
    for col, h in enumerate(headers):
        style_header_cell(table.cell(0, col), h)

    cum = 0
    for ri, entry in enumerate(ACTUAL_WEEKLY, start=1):
        wt = weekly_total(entry)
        cum += wt
        table.cell(ri, 0).text = entry["week"]
        table.cell(ri, 1).text = str(entry["lite"])
        table.cell(ri, 2).text = str(entry["pro"])
        table.cell(ri, 3).text = str(entry["enterprise"])
        table.cell(ri, 4).text = str(wt)
        for c in range(5):
            table.cell(ri, c).text_frame.paragraphs[0].font.size = Pt(11)

    add_textbox(slide, Inches(1.5), Inches(6.45), Inches(10), Inches(0.3),
                f"Grand cumulative total through {CURRENT_AS_OF}: {CURRENT_CUMULATIVE:,}",
                size=12, bold=True, color=COLOR_ACTUAL, align=PP_ALIGN.CENTER)

    add_nav_bar(slide, slides_map, "weekly", google_compatible)


def build_phases_slide(slide, slides_map, google_compatible=False):
    add_bg(slide)
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.55),
                "Migration Phases (Editable)", size=24, bold=True)

    headers = ["Phase", "Start", "End", "Days", "Reservations", "Accounts", "Res/Acct"]
    table_shape = slide.shapes.add_table(
        len(MIGRATION_PHASES) + 1, len(headers),
        Inches(0.8), Inches(1.2), Inches(11.7), Inches(3.2),
    )
    table = table_shape.table
    for col, h in enumerate(headers):
        style_header_cell(table.cell(0, col), h)
    for ri, phase in enumerate(MIGRATION_PHASES, start=1):
        vals = [
            phase["phase"], phase["start"], phase["end"], str(phase["days"]),
            f"{phase['total_reservations']:,}", f"{phase['accounts_in_phase']:,}",
            str(phase["avg_res_per_account"]),
        ]
        for ci, val in enumerate(vals):
            table.cell(ri, ci).text = val
            table.cell(ri, ci).text_frame.paragraphs[0].font.size = Pt(10)

    add_nav_bar(slide, slides_map, "phases", google_compatible)


def ensure_chart_png():
    if not os.path.exists(CHART_PNG):
        subprocess.run([sys.executable, "/workspace/create_actual_vs_planned_chart.py"], check=True)


def build_deck(output_path="/workspace/migration_cumulative_slide.pptx", google_compatible=False):
    ensure_chart_png()
    rows = build_chart_rows()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s_overview = prs.slides.add_slide(blank)
    s_chart = prs.slides.add_slide(blank)
    s_data = prs.slides.add_slide(blank)
    s_weekly = prs.slides.add_slide(blank)
    s_phases = prs.slides.add_slide(blank)

    slides_map = {
        "overview": s_overview,
        "chart": s_chart,
        "data": s_data,
        "weekly": s_weekly,
        "phases": s_phases,
    }

    build_overview_slide(s_overview, slides_map, rows, google_compatible)
    build_chart_slide(s_chart, slides_map, google_compatible)
    build_data_slide(s_data, slides_map, rows, google_compatible)
    build_weekly_slide(s_weekly, slides_map, google_compatible)
    build_phases_slide(s_phases, slides_map, google_compatible)

    prs.save(output_path)
    mode = "Google-compatible" if google_compatible else "Interactive"
    print(f"Saved Google Slides deck to {output_path} ({mode})")
    print("Slides:")
    print("  1. Overview — summary cards + recent cumulative totals")
    print("  2. Chart — Actual vs Planned line chart")
    print("  3. Cumulative Data — editable weekly + cumulative table")
    print("  4. Weekly Detail — editable segment breakdown")
    print("  5. Migration Phases — editable phase plan")
    if google_compatible:
        print("\nImport: Google Drive → Upload file → Right-click → Open with Google Slides")
    else:
        print("\nImport: Upload to Google Drive → Open with Google Slides")


if __name__ == "__main__":
    import sys
    compatible = "--google-compatible" in sys.argv
    path = "/workspace/migration_cumulative_slide_google.pptx" if compatible else "/workspace/migration_cumulative_slide.pptx"
    if not os.path.exists(CHART_PNG):
        subprocess.run([sys.executable, "/workspace/create_actual_vs_planned_chart.py"], check=False)
    build_deck(path, google_compatible=compatible)
