#!/usr/bin/env python3
"""Generate Room Migration Progress summary slide from sheet data."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Data from Room_Migration_Progress sheet
WEEKLY_DATA = [
    {"week": "5/25/2026", "lite": 160, "pro": 97, "enterprise": 0, "mid_market": 0},
    {"week": "6/1/2026", "lite": 111, "pro": 135, "enterprise": 0, "mid_market": 0},
    {"week": "6/8/2026", "lite": 428, "pro": 575, "enterprise": 9, "mid_market": 0},
    {"week": "6/15/2026", "lite": 840, "pro": 1810, "enterprise": 0, "mid_market": 181},
]

COLOR_LITE = RGBColor(135, 206, 235)      # light blue
COLOR_PRO = RGBColor(27, 79, 138)         # dark blue
COLOR_ENTERPRISE = RGBColor(30, 41, 59)   # black/navy
COLOR_MID = RGBColor(46, 158, 91)         # teal/green
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(209, 213, 219)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def week_total(row):
    return row["lite"] + row["pro"] + row["enterprise"] + row["mid_market"]


def segment_totals():
    return {
        "lite": sum(r["lite"] for r in WEEKLY_DATA),
        "pro": sum(r["pro"] for r in WEEKLY_DATA),
        "enterprise": sum(r["enterprise"] for r in WEEKLY_DATA),
        "mid_market": sum(r["mid_market"] for r in WEEKLY_DATA),
    }


def grand_total():
    return sum(week_total(r) for r in WEEKLY_DATA)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def build_summary_slide(output_path="/workspace/room_migration_summary_slide.pptx"):
    totals = segment_totals()
    total = grand_total()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.7),
                "Room Migration Progress — Summary", size=28, bold=True, align=PP_ALIGN.CENTER)

    # Grand total card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(1.35), Inches(4.9), Inches(1.1))
    set_fill(card, COLOR_WHITE)
    card.line.color.rgb = COLOR_BORDER
    card.adjustments[0] = 0.1
    add_textbox(slide, Inches(4.4), Inches(1.45), Inches(4.5), Inches(0.35),
                "Total Accounts Migrated", size=14, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.4), Inches(1.85), Inches(4.5), Inches(0.5),
                f"{total:,}", size=36, bold=True, align=PP_ALIGN.CENTER, color=COLOR_PRO)

    # Segment breakdown cards
    segments = [
        ("Lite", totals["lite"], COLOR_LITE),
        ("Pro", totals["pro"], COLOR_PRO),
        ("Enterprise", totals["enterprise"], COLOR_ENTERPRISE),
        ("Mid Market", totals["mid_market"], COLOR_MID),
    ]
    card_w = Inches(2.8)
    card_h = Inches(1.6)
    gap = Inches(0.35)
    start_x = (SLIDE_W - (card_w * 4 + gap * 3)) / 2

    for i, (name, count, color) in enumerate(segments):
        left = start_x + i * (card_w + gap)
        top = Inches(2.85)
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        set_fill(c, COLOR_WHITE)
        c.line.color.rgb = COLOR_BORDER
        c.adjustments[0] = 0.08
        add_textbox(slide, left + Inches(0.15), top + Inches(0.15), card_w - Inches(0.3), Inches(0.3),
                    name, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.55), card_w - Inches(0.3), Inches(0.45),
                    f"{count:,}", size=26, bold=True, align=PP_ALIGN.CENTER)
        pct = count / total * 100 if total else 0
        add_textbox(slide, left + Inches(0.15), top + Inches(1.15), card_w - Inches(0.3), Inches(0.3),
                    f"{pct:.1f}% of total", size=10, align=PP_ALIGN.CENTER, color=RGBColor(100, 116, 139))

    # Weekly totals table
    add_textbox(slide, Inches(0.8), Inches(4.75), Inches(5), Inches(0.35),
                "Weekly Totals", size=14, bold=True)

    rows = len(WEEKLY_DATA) + 1
    table_shape = slide.shapes.add_table(rows, 6, Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.6))
    table = table_shape.table
    headers = ["Week Start", "Lite", "Pro", "Enterprise", "Mid Market", "Weekly Total"]
    for col, h in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRO
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

    for row_idx, entry in enumerate(WEEKLY_DATA, start=1):
        wt = week_total(entry)
        values = [
            entry["week"],
            f"{entry['lite']:,}",
            f"{entry['pro']:,}",
            f"{entry['enterprise']:,}" if entry["enterprise"] else "0",
            f"{entry['mid_market']:,}" if entry["mid_market"] else "0",
            f"{wt:,}",
        ]
        for col, val in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            if col == 5:
                cell.text_frame.paragraphs[0].font.bold = True

    prs.save(output_path)
    print(f"Saved summary slide to {output_path}")


def print_summary():
    totals = segment_totals()
    total = grand_total()
    print("\n=== Room Migration Progress — Totals ===\n")
    print(f"Grand Total: {total:,} accounts\n")
    print("By Segment (cumulative):")
    for name, key in [("Lite", "lite"), ("Pro", "pro"), ("Enterprise", "enterprise"), ("Mid Market", "mid_market")]:
        count = totals[key]
        pct = count / total * 100
        print(f"  {name:12} {count:>6,}  ({pct:.1f}%)")
    print("\nWeekly Totals:")
    for entry in WEEKLY_DATA:
        wt = week_total(entry)
        print(f"  {entry['week']:10}  {wt:>6,}")
    print(f"\n  {'Total':10}  {total:>6,}")


if __name__ == "__main__":
    print_summary()
    build_summary_slide()
