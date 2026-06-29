#!/usr/bin/env python3
"""Create editable weekly migration schedule slide (PPTX for Google Slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from migration_schedule_data import FINAL_TARGET, WEEKLY_SCHEDULE

COLOR_LITE = RGBColor(27, 79, 138)
COLOR_SMB = RGBColor(46, 158, 91)
COLOR_MID = RGBColor(156, 163, 175)
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CARD_BORDER = RGBColor(209, 213, 219)
COLOR_GHOST = RGBColor(209, 213, 219)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height, text, size=12, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(shape, fill_color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.25)
    else:
        set_no_line(shape)
    shape.adjustments[0] = radius
    return shape


def add_timeline(slide):
    bar_left = Inches(1.0)
    bar_top = Inches(1.65)
    bar_width = Inches(11.3)
    segment_w = bar_width / len(WEEKLY_SCHEDULE)

    for i, entry in enumerate(WEEKLY_SCHEDULE):
        color = COLOR_SMB if entry["status"] == "actual" else COLOR_GHOST
        alpha = 1.0 if entry["status"] == "actual" else 0.5
        seg = add_rounded_rect(slide, bar_left + segment_w * i, bar_top, int(segment_w * 0.95), Inches(0.22), color, radius=0.5)
        if entry["status"] == "planned":
            seg.fill.transparency = 0.35

        dot_size = Inches(0.22)
        x = bar_left + segment_w * i + segment_w / 2 - dot_size / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, bar_top + Inches(0.03), dot_size, dot_size)
        set_fill(dot, COLOR_LITE if entry["status"] == "actual" else COLOR_MID)
        set_no_line(dot)

        add_textbox(slide, bar_left + segment_w * i, bar_top + Inches(0.32), segment_w, Inches(0.3),
                    entry["week"], size=9, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, bar_left + segment_w * i, bar_top + Inches(0.55), segment_w, Inches(0.45),
                    entry["milestone"], size=7, align=PP_ALIGN.CENTER, italic=entry["status"] == "planned")


def add_progress_summary(slide):
    actual_weeks = [w for w in WEEKLY_SCHEDULE if w["status"] == "actual"]
    current = actual_weeks[-1]
    prev = actual_weeks[-2] if len(actual_weeks) > 1 else None
    pct_to_goal = current["target_total"] / FINAL_TARGET * 100

    summary_top = Inches(2.35)
    card_w = Inches(3.8)
    card_h = Inches(1.5)
    gap = Inches(0.35)
    start_left = (SLIDE_W - (card_w * 3 + gap * 2)) / 2

    cards = [
        ("Current Total", f"{current['target_total']:,}", f"As of {current['week']}", COLOR_SMB),
        ("Week-over-Week", f"+{((current['target_total'] - prev['target_total']) / prev['target_total'] * 100):.1f}%" if prev else "—",
         f"{prev['week']} → {current['week']}" if prev else "", COLOR_LITE),
        ("Progress to Goal", f"{pct_to_goal:.1f}%", f"Target: {FINAL_TARGET:,}", COLOR_MID),
    ]

    for i, (title, value, subtitle, color) in enumerate(cards):
        left = start_left + i * (card_w + gap)
        card = add_rounded_rect(slide, left, summary_top, card_w, card_h, COLOR_WHITE, COLOR_CARD_BORDER)
        add_textbox(slide, left + Inches(0.2), summary_top + Inches(0.15), card_w - Inches(0.4), Inches(0.3), title, size=11, color=color, bold=True)
        add_textbox(slide, left + Inches(0.2), summary_top + Inches(0.5), card_w - Inches(0.4), Inches(0.45), value, size=22, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.2), summary_top + Inches(1.05), card_w - Inches(0.4), Inches(0.3), subtitle, size=9, color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)


def add_schedule_table_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
                "Weekly Schedule Data (Editable)", size=24, bold=True)

    rows, cols = len(WEEKLY_SCHEDULE) + 1, 8
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5))
    table = table_shape.table

    headers = ["Week", "Label", "Milestone", "Total", "LITE", "SMB", "Mid-Market", "Status"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LITE
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

    for row_idx, entry in enumerate(WEEKLY_SCHEDULE, start=1):
        values = [
            entry["week"],
            entry["week_label"],
            entry["milestone"],
            str(entry["target_total"]),
            str(entry["target_lite"]),
            str(entry["target_smb"]),
            str(entry["target_mid"]),
            entry["status"],
        ]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1),
                f"Edit the table above to update the chart. Final target: {FINAL_TARGET:,} accounts.",
                size=11, color=COLOR_MID, italic=True)


def build_slide(output_path="/workspace/migration_weekly_schedule_slide.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.75),
                "Reservation Migration Plan — Weekly Progress", size=28, bold=True, align=PP_ALIGN.CENTER)

    add_timeline(slide)
    add_progress_summary(slide)

    # Mini stacked bars for each week
    bar_area_top = Inches(4.1)
    bar_h = Inches(2.2)
    bar_w = Inches(1.7)
    gap = Inches(0.25)
    total_w = bar_w * len(WEEKLY_SCHEDULE) + gap * (len(WEEKLY_SCHEDULE) - 1)
    bar_left_start = (SLIDE_W - total_w) / 2
    max_total = max(w["target_total"] for w in WEEKLY_SCHEDULE)

    for i, entry in enumerate(WEEKLY_SCHEDULE):
        left = bar_left_start + i * (bar_w + gap)
        total = entry["target_total"]
        scale = bar_h / max_total

        lite_h = int(entry["target_lite"] * scale)
        smb_h = int(entry["target_smb"] * scale)
        mid_h = int(entry["target_mid"] * scale)

        y = bar_area_top + bar_h
        if mid_h:
            add_rounded_rect(slide, left, y - mid_h, bar_w, mid_h, COLOR_MID, radius=0.15)
            y -= mid_h
        if smb_h:
            add_rounded_rect(slide, left, y - smb_h, bar_w, smb_h, COLOR_SMB, radius=0.15)
            y -= smb_h
        if lite_h:
            add_rounded_rect(slide, left, y - lite_h, bar_w, lite_h, COLOR_LITE, radius=0.15)

        add_textbox(slide, left, bar_area_top + bar_h + Inches(0.08), bar_w, Inches(0.25),
                    f"{total:,}", size=9, bold=True, align=PP_ALIGN.CENTER)
        status = "✓" if entry["status"] == "actual" else "○"
        add_textbox(slide, left, bar_area_top - Inches(0.25), bar_w, Inches(0.2),
                    status, size=10, align=PP_ALIGN.CENTER, color=COLOR_SMB if entry["status"] == "actual" else COLOR_MID)

    add_textbox(slide, Inches(0.8), Inches(6.85), Inches(2), Inches(0.4),
                "Guesty", size=18, bold=True, color=COLOR_LITE, italic=True)
    add_textbox(slide, Inches(9.5), Inches(6.85), Inches(3), Inches(0.4),
                "○ Planned   ✓ Actual", size=10, color=COLOR_MID, align=PP_ALIGN.RIGHT)

    add_schedule_table_slide(prs)
    prs.save(output_path)
    print(f"Saved editable slide to {output_path}")
    print("Slides included:")
    print("  1. Visual weekly progress (timeline + summary + bars)")
    print("  2. Editable schedule data table")


if __name__ == "__main__":
    build_slide()
