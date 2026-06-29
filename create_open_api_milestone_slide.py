#!/usr/bin/env python3
"""Editable Google Slides deck: Open API milestone timeline bar."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Milestone data (edit to update the timeline) ---
TIMELINE_START = "Jul 2026"
TIMELINE_END = "Mar 2027"

EXISTING_USERS_MILESTONES = [
    {
        "date": "Sep 2026",
        "date_detail": "Beginning of Sept",
        "title": "Automated email to migrated users",
        "description": "Start sending automated email to migrated users",
    },
    {
        "date": "Sep 2026",
        "date_detail": "End of September",
        "title": "Communication to all users",
        "description": "Sending communication to all users",
    },
    {
        "date": "Mar 27, 2027",
        "date_detail": "EO March 27",
        "title": "Deprecation of old Open API",
        "description": "End of deprecation period for old OAPI",
    },
]

NEW_USERS_MILESTONES = [
    {
        "date": "Early Q3 2026",
        "date_detail": "Early Q3",
        "title": "Block API usage by registration date",
        "description": "By registration date, block API usage for new users",
    },
]

# Position on timeline bar (0.0 = left / start, 1.0 = right / end)
MILESTONE_POSITIONS = {
    "Early Q3 2026": 0.05,
    "Sep 2026 (begin)": 0.42,
    "Sep 2026 (end)": 0.52,
    "Mar 27, 2027": 0.95,
}

COLOR_EXISTING = RGBColor(27, 79, 138)    # dark blue
COLOR_NEW = RGBColor(46, 158, 91)          # green
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BAR = RGBColor(209, 213, 219)
COLOR_BAR_FILL = RGBColor(27, 79, 138)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)


def add_textbox(slide, left, top, width, height, text, size=12, bold=False,
                color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_timeline_bar(slide, track_top, track_label, milestones, color, positions):
    """Draw a labeled timeline track with milestone markers."""
    bar_left = Inches(1.2)
    bar_width = Inches(10.9)
    bar_height = Inches(0.12)

    # Track label
    add_textbox(slide, Inches(0.3), track_top - Inches(0.05), Inches(0.85), Inches(0.9),
                track_label, size=9, bold=True, color=color, align=PP_ALIGN.RIGHT)

    # Background bar
    add_rounded_rect(slide, bar_left, track_top + Inches(0.35), bar_width, bar_height,
                     COLOR_BAR, radius=0.5)

    # Filled portion (progress to last milestone on this track)
    if milestones:
        last_pos = max(positions.get(m["position_key"], 0.5) for m in milestones)
        fill_w = int(bar_width * last_pos)
        if fill_w > 0:
            add_rounded_rect(slide, bar_left, track_top + Inches(0.35), fill_w, bar_height,
                             color, radius=0.5)

    for m in milestones:
        pos = positions.get(m["position_key"], 0.5)
        x = bar_left + int(bar_width * pos)
        dot_size = Inches(0.22)
        dot_y = track_top + Inches(0.3)

        # Marker dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - dot_size // 2, dot_y, dot_size, dot_size)
        set_fill(dot, color)
        set_no_line(dot)

        # Date label above bar
        add_textbox(slide, x - Inches(0.65), track_top + Inches(0.02), Inches(1.3), Inches(0.28),
                    m["date_detail"], size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

        # Milestone card below bar
        card_w = Inches(2.6)
        card_h = Inches(0.95)
        card_left = x - card_w // 2
        card_top = track_top + Inches(0.55)

        card = add_rounded_rect(slide, card_left, card_top, card_w, card_h, COLOR_WHITE, color)
        card.line.width = Pt(1.5)

        add_textbox(slide, card_left + Inches(0.12), card_top + Inches(0.08),
                    card_w - Inches(0.24), Inches(0.35),
                    m["title"], size=9, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, card_left + Inches(0.12), card_top + Inches(0.42),
                    card_w - Inches(0.24), Inches(0.48),
                    m["description"], size=8, color=COLOR_TEXT, align=PP_ALIGN.CENTER, italic=True)


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(shape, fill_color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.25)
    else:
        set_no_line(shape)
    shape.adjustments[0] = radius
    return shape


def build_timeline_slide(slide):
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.65),
                "Open API — Milestone Timeline", size=28, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11.5), Inches(0.35),
                f"{TIMELINE_START} → {TIMELINE_END}  ·  Click any text box or table cell to edit",
                size=11, color=RGBColor(100, 116, 139), align=PP_ALIGN.CENTER, italic=True)

    # Timeline axis labels
    bar_left = Inches(1.2)
    bar_width = Inches(10.9)
    axis_y = Inches(1.55)
    add_textbox(slide, bar_left, axis_y, Inches(1.2), Inches(0.3), TIMELINE_START,
                size=10, bold=True, color=COLOR_TEXT)
    add_textbox(slide, bar_left + bar_width - Inches(1.2), axis_y, Inches(1.2), Inches(0.3),
                TIMELINE_END, size=10, bold=True, color=COLOR_TEXT, align=PP_ALIGN.RIGHT)

    # Prepare milestone dicts with position keys
    existing = []
    for i, m in enumerate(EXISTING_USERS_MILESTONES):
        key = "Sep 2026 (begin)" if i == 0 else "Sep 2026 (end)" if i == 1 else "Mar 27, 2027"
        existing.append({**m, "position_key": key})

    new = [{**NEW_USERS_MILESTONES[0], "position_key": "Early Q3 2026"}]

    add_timeline_bar(slide, Inches(2.0), "New\nUsers", new, COLOR_NEW, MILESTONE_POSITIONS)
    add_timeline_bar(slide, Inches(4.1), "Existing\nUsers", existing, COLOR_EXISTING, MILESTONE_POSITIONS)

    # Legend
    legend_y = Inches(6.55)
    for i, (label, color) in enumerate([("New Users", COLOR_NEW), ("Existing Users", COLOR_EXISTING)]):
        left = Inches(4.5) + i * Inches(2.5)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, legend_y, Inches(0.2), Inches(0.2))
        set_fill(sq, color)
        set_no_line(sq)
        add_textbox(slide, left + Inches(0.3), legend_y - Inches(0.05), Inches(2), Inches(0.3),
                    label, size=10, color=COLOR_TEXT)

    add_textbox(slide, Inches(0.8), Inches(6.85), Inches(2), Inches(0.35),
                "Guesty", size=14, bold=True, color=COLOR_EXISTING, italic=True)


def build_data_slide(slide):
    """Editable milestone table for easy updates."""
    add_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.55),
                "Milestone Data (Editable)", size=24, bold=True)
    add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35),
                "Edit this table, then update the timeline slide to match.",
                size=11, color=RGBColor(100, 116, 139), italic=True)

    header_fmt_bg = COLOR_EXISTING
    rows_data = []

    rows_data.append(["Segment", "Date", "Date Label", "Milestone Title", "Description"])
    for m in NEW_USERS_MILESTONES:
        rows_data.append(["New Users", m["date"], m["date_detail"], m["title"], m["description"]])
    for m in EXISTING_USERS_MILESTONES:
        rows_data.append(["Existing Users", m["date"], m["date_detail"], m["title"], m["description"]])

    cols = 5
    table_shape = slide.shapes.add_table(len(rows_data), cols, Inches(0.6), Inches(1.5),
                                         Inches(12.1), Inches(2.8))
    table = table_shape.table
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.6)
    table.columns[3].width = Inches(3.5)
    table.columns[4].width = Inches(4.1)

    for ri, row in enumerate(rows_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10 if ri > 0 else 11)
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fmt_bg
            elif row[0] == "New Users":
                if ci == 0:
                    p.font.color.rgb = COLOR_NEW
                    p.font.bold = True

    add_textbox(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2),
                "Timeline span: Jul 2026 → Mar 2027\n\n"
                "New Users\n"
                "  • Early Q3 — Block API usage by registration date\n\n"
                "Existing Users\n"
                "  • Beginning of Sept — Automated email to migrated users\n"
                "  • End of September — Communication to all users\n"
                "  • EO March 27 — Deprecation of old Open API (OAPI)",
                size=12, color=COLOR_TEXT)


def build_deck(output_path="/workspace/open_api_milestone_timeline.pptx", google_compatible=False):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s_timeline = prs.slides.add_slide(blank)
    s_data = prs.slides.add_slide(blank)

    build_timeline_slide(s_timeline)
    build_data_slide(s_data)

    if not google_compatible:
        # Internal slide links break Google Slides import — skip in compatible mode
        hint = s_timeline.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(10.5), Inches(6.85), Inches(2.3), Inches(0.4))
        set_fill(hint, COLOR_NEW)
        set_no_line(hint)
        hint.adjustments[0] = 0.4
        hint.click_action.target_slide = s_data
        tf = hint.text_frame
        tf.paragraphs[0].text = "✎ Edit milestone data →"
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        add_textbox(s_timeline, Inches(10.2), Inches(6.85), Inches(2.6), Inches(0.4),
                    "Slide 2: Milestone Data", size=10, bold=True, color=COLOR_NEW, align=PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"Saved to {output_path}")
    mode = "Google-compatible" if google_compatible else "Interactive"
    print(f"Mode: {mode}")
    print("Slides:")
    print("  1. Milestone Timeline — visual timeline bar (editable text boxes)")
    print("  2. Milestone Data — editable table")
    if google_compatible:
        print("\nImport: Google Drive → Upload → Right-click → Open with Google Slides")


if __name__ == "__main__":
    import sys
    compatible = "--google-compatible" in sys.argv
    path = "/workspace/open_api_milestone_timeline_google.pptx" if compatible else "/workspace/open_api_milestone_timeline.pptx"
    build_deck(path, google_compatible=compatible)
