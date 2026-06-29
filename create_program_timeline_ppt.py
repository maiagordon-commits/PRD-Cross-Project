#!/usr/bin/env python3
"""Generate import-friendly PowerPoint program timeline (Google Slides / PowerPoint).

No internal hyperlinks, standard Arial font, simple shapes only.
"""

from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from timeline_data import (
    OPEN_API_MILESTONES,
    OPEN_API_TIMELINE_START,
    OPEN_API_TIMELINE_END,
    build_master_timeline,
    build_reservation_phase_rows,
    parse_iso,
)

# Guesty palette
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_MID = RGBColor(100, 116, 139)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_OPEN_API = RGBColor(27, 79, 138)
COLOR_NEW = RGBColor(46, 158, 91)
COLOR_MIGRATION = RGBColor(217, 119, 6)
COLOR_RESERVATION = RGBColor(124, 58, 237)
COLOR_BAR = RGBColor(209, 213, 219)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Arial"

TIMELINE_START = parse_iso("2026-06-01")
TIMELINE_END = parse_iso("2027-03-31")
BAR_LEFT = Inches(1.35)
BAR_WIDTH = Inches(10.75)


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def style_paragraph(p, size=12, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False):
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.italic = italic
    p.alignment = align


def add_textbox(slide, left, top, width, height, text, size=12, bold=False,
                color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    style_paragraph(tf.paragraphs[0], size, bold, color, align, italic)
    tf.paragraphs[0].text = text
    return box


def date_to_fraction(date_str):
    dt = parse_iso(date_str)
    total = (TIMELINE_END - TIMELINE_START).days
    if total <= 0:
        return 0.0
    return max(0.0, min(1.0, (dt - TIMELINE_START).days / total))


def add_axis(slide, top):
    add_textbox(slide, BAR_LEFT, top, Inches(1.2), Inches(0.28), "Jun 2026",
                size=9, bold=True, color=COLOR_MID)
    add_textbox(slide, BAR_LEFT + BAR_WIDTH - Inches(1.4), top, Inches(1.4), Inches(0.28),
                "Mar 2027", size=9, bold=True, color=COLOR_MID, align=PP_ALIGN.RIGHT)
    mid = BAR_LEFT + BAR_WIDTH // 2 - Inches(0.5)
    add_textbox(slide, mid, top, Inches(1.0), Inches(0.28), "Jan 2027",
                size=9, color=COLOR_MID, align=PP_ALIGN.CENTER)


def draw_track_bar(slide, top, label, color, start_date, end_date):
    """Horizontal duration bar between two dates."""
    x0 = BAR_LEFT + int(BAR_WIDTH * date_to_fraction(start_date))
    x1 = BAR_LEFT + int(BAR_WIDTH * date_to_fraction(end_date))
    width = max(Inches(0.08), x1 - x0)
    height = Inches(0.14)

    add_textbox(slide, Inches(0.15), top, Inches(1.1), Inches(0.55),
                label, size=8, bold=True, color=color, align=PP_ALIGN.RIGHT)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BAR_LEFT, top + Inches(0.38), BAR_WIDTH, height)
    set_fill(bg, COLOR_BAR)
    set_no_line(bg)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, top + Inches(0.38), width, height)
    set_fill(bar, color)
    set_no_line(bar)


def draw_point_milestone(slide, top, date_str, label, title, color):
    """Single-date milestone marker on a track."""
    x = BAR_LEFT + int(BAR_WIDTH * date_to_fraction(date_str))
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.09), top + Inches(0.32), Inches(0.18), Inches(0.18))
    set_fill(dot, color)
    set_no_line(dot)
    add_textbox(slide, x - Inches(0.55), top + Inches(0.02), Inches(1.1), Inches(0.25),
                label, size=8, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, x - Inches(1.0), top + Inches(0.55), Inches(2.0), Inches(0.55),
                title, size=8, color=COLOR_TEXT, align=PP_ALIGN.CENTER)


def api_date_to_fraction(date_str):
    """Position on Open API timeline (Jul 2026 → Mar 2027)."""
    start = parse_iso(OPEN_API_TIMELINE_START)
    end = parse_iso(OPEN_API_TIMELINE_END)
    dt = parse_iso(date_str)
    total = (end - start).days
    if total <= 0:
        return 0.0
    return max(0.0, min(1.0, (dt - start).days / total))


def draw_unified_open_api_bar(slide, track_top):
    """Single timeline bar: green (New Users) + blue (Existing Users) with all milestones."""
    bar_h = Inches(0.18)
    bar_y = track_top + Inches(0.55)

    # Gray background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BAR_LEFT, bar_y, BAR_WIDTH, bar_h)
    set_fill(bg, COLOR_BAR)
    set_no_line(bg)

    # Two-tone fill: green until Existing Users phase, then blue through end
    split_frac = api_date_to_fraction("2026-09-01")
    green_w = int(BAR_WIDTH * split_frac)
    if green_w > 0:
        green = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, BAR_LEFT, bar_y, green_w, bar_h)
        set_fill(green, COLOR_NEW)
        set_no_line(green)
    blue_left = BAR_LEFT + green_w
    blue_w = BAR_WIDTH - green_w
    if blue_w > 0:
        blue = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, blue_left, bar_y, blue_w, bar_h)
        set_fill(blue, COLOR_OPEN_API)
        set_no_line(blue)

    add_textbox(slide, Inches(0.15), track_top + Inches(0.35), Inches(1.1), Inches(0.45),
                "Open API\nMilestones", size=9, bold=True, color=COLOR_TEXT, align=PP_ALIGN.RIGHT)

    # Milestone markers on the single bar (cards alternate above / below)
    for i, m in enumerate(OPEN_API_MILESTONES):
        is_new = m["track"] == "New Users"
        color = COLOR_NEW if is_new else COLOR_OPEN_API
        x = BAR_LEFT + int(BAR_WIDTH * api_date_to_fraction(m["start_date"]))
        dot_y = bar_y - Inches(0.04)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.1), dot_y, Inches(0.2), Inches(0.2))
        set_fill(dot, color)
        set_no_line(dot)
        dot.line.color.rgb = COLOR_WHITE
        dot.line.width = Pt(2)

        above = i % 2 == 0
        if above:
            date_top = track_top - Inches(0.05)
            card_top = track_top + Inches(0.05)
        else:
            date_top = bar_y + Inches(0.35)
            card_top = bar_y + Inches(0.55)

        add_textbox(slide, x - Inches(0.65), date_top, Inches(1.3), Inches(0.28),
                    m["date_label"], size=9, bold=True, color=color, align=PP_ALIGN.CENTER)

        card_w = Inches(2.55)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x - card_w // 2, card_top, card_w, Inches(0.95))
        set_fill(card, COLOR_WHITE)
        card.line.color.rgb = color
        card.line.width = Pt(1.5)
        card.adjustments[0] = 0.06
        seg = "New Users" if is_new else "Existing Users"
        add_textbox(slide, x - card_w // 2 + Inches(0.1), card_top + Inches(0.06),
                    card_w - Inches(0.2), Inches(0.22), seg, size=7, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - card_w // 2 + Inches(0.1), card_top + Inches(0.28),
                    card_w - Inches(0.2), Inches(0.35), m["milestone"], size=8, bold=True,
                    color=COLOR_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, x - card_w // 2 + Inches(0.1), card_top + Inches(0.62),
                    card_w - Inches(0.2), Inches(0.3), m["description"], size=7,
                    color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)


def build_open_api_slide(slide):
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.55),
                "Open API — Milestone Timeline", size=26, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(0.88), Inches(12), Inches(0.3),
                "Jul 2026 → Mar 2027 · Green = New Users · Blue = Existing Users",
                size=11, color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)

    add_textbox(slide, BAR_LEFT, Inches(1.28), Inches(1.2), Inches(0.28), "Jul 2026",
                size=9, bold=True, color=COLOR_MID)
    add_textbox(slide, BAR_LEFT + BAR_WIDTH - Inches(1.4), Inches(1.28), Inches(1.4), Inches(0.28),
                "Mar 2027", size=9, bold=True, color=COLOR_MID, align=PP_ALIGN.RIGHT)

    draw_unified_open_api_bar(slide, Inches(1.55))

    # Legend
    legend_y = Inches(6.35)
    for i, (label, color) in enumerate([("New Users", COLOR_NEW), ("Existing Users", COLOR_OPEN_API)]):
        left = Inches(4.8) + i * Inches(2.8)
        sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, legend_y, Inches(0.22), Inches(0.22))
        set_fill(sq, color)
        set_no_line(sq)
        add_textbox(slide, left + Inches(0.32), legend_y - Inches(0.04), Inches(2.2), Inches(0.3),
                    label, size=10, color=COLOR_TEXT)

    add_textbox(slide, Inches(0.6), Inches(6.75), Inches(2), Inches(0.3),
                "Guesty", size=13, bold=True, color=COLOR_OPEN_API, italic=True)


def build_migration_phases_slide(slide):
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.55),
                "Reservation Migration — Phase Timeline", size=26, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(0.88), Inches(12), Inches(0.3),
                "Jun 15 – Dec 31, 2026 · ~35K accounts · ~54M reservations",
                size=11, color=COLOR_MID, align=PP_ALIGN.CENTER, italic=True)
    add_axis(slide, Inches(1.25))

    phases = build_reservation_phase_rows()
    colors = [COLOR_NEW, COLOR_OPEN_API, COLOR_MIGRATION, COLOR_RESERVATION]
    top = Inches(1.65)
    step = Inches(1.05)

    for i, phase in enumerate(phases):
        y = top + i * step
        color = colors[i % len(colors)]
        draw_track_bar(slide, y, phase["track"], color, phase["start_date"], phase["end_date"])
        x_end = BAR_LEFT + int(BAR_WIDTH * date_to_fraction(phase["end_date"]))
        add_textbox(slide, x_end + Inches(0.08), y + Inches(0.28), Inches(2.2), Inches(0.45),
                    f"{phase['target_accounts']:,} accts", size=8, bold=True, color=color)
        add_textbox(slide, BAR_LEFT, y + Inches(0.55), BAR_WIDTH, Inches(0.35),
                    phase["description"], size=8, color=COLOR_MID, italic=True)

    add_textbox(slide, Inches(0.6), Inches(6.75), Inches(2), Inches(0.3),
                "Guesty", size=13, bold=True, color=COLOR_OPEN_API, italic=True)


def build_data_slide(slide):
    set_slide_bg(slide)
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.5),
                "Full Program Timeline — Data (Editable)", size=24, bold=True)
    add_textbox(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3),
                "Edit cells below, then update the visual slides to match.",
                size=11, color=COLOR_MID, italic=True)

    rows = build_master_timeline()
    headers = ["Program", "Track", "Date", "Milestone", "Description", "Status"]
    table_data = [headers]
    for r in rows:
        table_data.append([
            r["program"], r["track"], r["date_label"],
            r["milestone"], r["description"], r["status"],
        ])

    shape = slide.shapes.add_table(len(table_data), len(headers),
                                   Inches(0.45), Inches(1.25), Inches(12.4), Inches(5.5))
    table = shape.table
    col_widths = [Inches(1.5), Inches(1.3), Inches(1.2), Inches(2.8), Inches(4.2), Inches(0.9)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ri, row in enumerate(table_data):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            p.font.size = Pt(9 if ri > 0 else 10)
            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_OPEN_API
            else:
                p.font.color.rgb = COLOR_TEXT


def build_deck(output_path="/workspace/guesty_program_timeline.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    s2 = prs.slides.add_slide(blank)
    s3 = prs.slides.add_slide(blank)

    build_open_api_slide(s1)
    build_migration_phases_slide(s2)
    build_data_slide(s3)

    prs.save(output_path)
    print(f"Saved: {output_path}")
    print("Slides:")
    print("  1. Open API Milestone Timeline")
    print("  2. Reservation Migration Phase Timeline")
    print("  3. Full Program Data (editable table)")
    print("\nImport into Google Slides:")
    print("  1. Upload .pptx to Google Drive")
    print("  2. Right-click → Open with → Google Slides")


if __name__ == "__main__":
    build_deck()
