#!/usr/bin/env python3
"""Create an editable migration progress slide (PPTX for Google Slides import)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Data (edit these values to update the slide) ---
MIGRATION_DATA = [
    {
        "date": "June 8th",
        "total": 1097,
        "lite": 380,
        "smb": 717,
        "mid_market": "0 (0.0%)",
        "pct_change": None,
    },
    {
        "date": "June 15th",
        "total": 2650,
        "lite": 840,
        "smb": 1810,
        "mid_market": "Pending",
        "pct_change": 141.6,
    },
]

# Guesty-style palette
COLOR_LITE = RGBColor(27, 79, 138)
COLOR_SMB = RGBColor(46, 158, 91)
COLOR_BG = RGBColor(245, 240, 232)
COLOR_TEXT = RGBColor(30, 41, 59)
COLOR_MID = RGBColor(156, 163, 175)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_CARD_BORDER = RGBColor(209, 213, 219)
COLOR_GHOST = RGBColor(209, 213, 219)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_no_line(shape):
    shape.line.fill.background()


def add_textbox(
    slide, left, top, width, height, text,
    size=12, bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT, italic=False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, radius=0.08):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(shape, fill_color)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.25)
    else:
        set_no_line(shape)
    shape.adjustments[0] = radius
    return shape


def add_mini_stacked_bar(slide, left, top, width, height, lite_count, total):
    lite_frac = lite_count / total
    lite_w = int(width * lite_frac)
    smb_w = width - lite_w

    if lite_w > 0:
        add_rounded_rect(slide, left, top, lite_w, height, COLOR_LITE, radius=0.2)
    if smb_w > 0:
        add_rounded_rect(slide, left + lite_w, top, smb_w, height, COLOR_SMB, radius=0.2)


def add_data_card(slide, left, top, width, height, entry):
    total = entry["total"]
    lite = entry["lite"]
    smb = entry["smb"]
    lite_pct = lite / total * 100
    smb_pct = smb / total * 100

    card = add_rounded_rect(slide, left, top, width, height, COLOR_WHITE, COLOR_CARD_BORDER)

    pad = Inches(0.35)
    inner_w = width - pad * 2
    y = top + Inches(0.3)

    # Date header
    add_textbox(slide, left + pad, y, inner_w, Inches(0.45), entry["date"], size=20, bold=True)
    y += Inches(0.55)

    # Growth badge (June 15th only)
    if entry["pct_change"] is not None:
        badge_w = Inches(1.1)
        badge = add_rounded_rect(slide, left + pad, y, badge_w, Inches(0.35), COLOR_SMB, radius=0.3)
        badge_tf = badge.text_frame
        badge_tf.paragraphs[0].text = f"+{entry['pct_change']:.1f}%"
        badge_tf.paragraphs[0].font.size = Pt(12)
        badge_tf.paragraphs[0].font.bold = True
        badge_tf.paragraphs[0].font.color.rgb = COLOR_WHITE
        badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge_tf.vertical_anchor = 1
        y += Inches(0.5)

    # Total accounts - large number
    add_textbox(
        slide, left + pad, y, inner_w, Inches(0.35),
        f"Total Accounts: {total:,}", size=16, bold=True,
    )
    y += Inches(0.55)

    # Segment rows
    add_textbox(
        slide, left + pad, y, inner_w, Inches(0.32),
        f"LITE Segment: {lite:,} ({lite_pct:.1f}%)", size=14, color=COLOR_LITE, bold=True,
    )
    y += Inches(0.42)

    add_textbox(
        slide, left + pad, y, inner_w, Inches(0.32),
        f"SMB Segment: {smb:,} ({smb_pct:.1f}%)", size=14, color=COLOR_SMB, bold=True,
    )
    y += Inches(0.42)

    add_textbox(
        slide, left + pad, y, inner_w, Inches(0.32),
        f"Mid-Market: {entry['mid_market']}", size=14, color=COLOR_MID, italic=True,
    )
    y += Inches(0.55)

    # Mini stacked bar
    add_mini_stacked_bar(slide, left + pad, y, inner_w, Inches(0.22), lite, total)


def add_timeline(slide):
    bar_left = Inches(1.5)
    bar_top = Inches(1.75)
    bar_width = Inches(10.3)
    bar_height = Inches(0.28)

    # Progress segments: blue -> green -> ghost
    blue_w = int(bar_width * 0.35)
    green_w = int(bar_width * 0.25)
    ghost_w = bar_width - blue_w - green_w

    add_rounded_rect(slide, bar_left, bar_top, blue_w, bar_height, COLOR_LITE, radius=0.5)
    add_rounded_rect(slide, bar_left + blue_w, bar_top, green_w, bar_height, COLOR_SMB, radius=0.5)

    ghost = add_rounded_rect(slide, bar_left + blue_w + green_w, bar_top, ghost_w, bar_height, COLOR_GHOST, radius=0.5)
    ghost.fill.transparency = 0.35

    # Timeline dots and labels
    markers = [
        (0.04, COLOR_LITE, "June 8th"),
        (0.58, COLOR_SMB, "June 15th"),
        (0.92, COLOR_SMB, ""),  # future milestone
    ]
    dot_size = Inches(0.28)
    for x_frac, color, label in markers:
        x = bar_left + int(bar_width * x_frac)
        dot_y = bar_top + bar_height / 2 - dot_size / 2
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, dot_y, dot_size, dot_size)
        set_fill(dot, color)
        set_no_line(dot)
        if label:
            add_textbox(
                slide, x - Inches(0.55), bar_top + Inches(0.42), Inches(1.3), Inches(0.35),
                label, size=11, align=PP_ALIGN.CENTER,
            )


def add_data_table_slide(prs):
    """Second slide with raw editable data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    add_textbox(
        slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
        "Migration Data (Editable)", size=24, bold=True,
    )

    rows, cols = 3, 6
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.8))
    table = table_shape.table

    headers = ["Date", "Total Accounts", "LITE", "SMB", "Mid-Market", "Change %"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LITE
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

    for row_idx, entry in enumerate(MIGRATION_DATA, start=1):
        values = [
            entry["date"],
            f"{entry['total']:,}",
            f"{entry['lite']:,} ({entry['lite'] / entry['total'] * 100:.1f}%)",
            f"{entry['smb']:,} ({entry['smb'] / entry['total'] * 100:.1f}%)",
            entry["mid_market"],
            f"+{entry['pct_change']:.1f}%" if entry["pct_change"] else "—",
        ]
        for col, value in enumerate(values):
            cell = table.cell(row_idx, col)
            cell.text = value
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    add_textbox(
        slide, Inches(0.8), Inches(3.6), Inches(11), Inches(2),
        "Edit the table above, then update the visual slide to match.\n"
        "Tip: In Google Slides, click any text box or table cell to edit directly.",
        size=12, color=COLOR_MID, italic=True,
    )


def build_slide(output_path="/workspace/migration_progress_slide.pptx"):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    set_fill(bg, COLOR_BG)
    set_no_line(bg)
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)

    # Title
    add_textbox(
        slide, Inches(0.8), Inches(0.55), Inches(11.5), Inches(0.75),
        "Existing Accounts Migration Progress", size=30, bold=True, align=PP_ALIGN.CENTER,
    )

    add_timeline(slide)

    # Two data cards side by side (matching original screenshot layout)
    card_w = Inches(5.2)
    card_h = Inches(3.6)
    card_top = Inches(2.65)
    gap = Inches(0.6)
    total_cards_w = card_w * 2 + gap
    card_left_start = (SLIDE_W - total_cards_w) / 2

    add_data_card(slide, card_left_start, card_top, card_w, card_h, MIGRATION_DATA[0])
    add_data_card(slide, card_left_start + card_w + gap, card_top, card_w, card_h, MIGRATION_DATA[1])

    # Branding
    add_textbox(
        slide, Inches(0.8), Inches(6.85), Inches(2), Inches(0.4),
        "Guesty", size=18, bold=True, color=COLOR_LITE, italic=True,
    )

    # Data table slide for easy editing
    add_data_table_slide(prs)

    prs.save(output_path)
    print(f"Saved editable slide to {output_path}")
    print("Slides included:")
    print("  1. Visual progress slide (timeline + data cards)")
    print("  2. Editable data table")
    print("Upload to Google Drive, then open with Google Slides to edit.")


if __name__ == "__main__":
    build_slide()
